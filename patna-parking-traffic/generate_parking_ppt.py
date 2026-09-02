#!/usr/bin/env python3
"""Patna Parking & Traffic Management presentation (15 slides).

Session by Pourush Agarwal, DySP Traffic Patna.
Uniform font sizes across all slides; real Patna photographs.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image as PILImage
import os

# ---------- Theme ----------
NAVY = RGBColor(0x0B, 0x2E, 0x4F)
BLUE = RGBColor(0x1B, 0x6C, 0xA8)
ORANGE = RGBColor(0xF2, 0x8C, 0x28)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6C, 0x3A, 0x8E)
PALE = RGBColor(0xAF, 0xC3, 0xD6)

# ---------- Uniform font sizes (same across every slide) ----------
F_TITLE = 27     # slide headers
F_BODY = 19      # all bullet text
F_CARD_H = 19    # card headings
F_CARD_T = 15    # card body text
F_CAP = 12       # photo captions

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ppt_images")


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, color, line=None, lw=1.5):
    sp = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def bullets(slide, l, t, w, h, items, size=F_BODY, color=NAVY, gap=9):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, sub = item if isinstance(item, tuple) else (item, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("•  " if not sub else "     – ") + text
        r.font.size = Pt(size if not sub else size - 3)
        r.font.color.rgb = color if not sub else GREY
        r.font.name = "Calibri"
    return tb


def header(slide, title, num):
    box(slide, 0, 0, SW, Inches(1.15), NAVY)
    box(slide, 0, Inches(1.15), SW, Inches(0.08), ORANGE)
    txt(slide, Inches(0.5), Inches(0.2), Inches(11.4), Inches(0.8), title,
        F_TITLE, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, Inches(12.35), Inches(0.2), Inches(0.7), Inches(0.8), str(num),
        15, ORANGE, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def cover_pic(slide, l, t, w, h, filename):
    """Place a picture that fills the box exactly, cropping (never stretching)."""
    path = os.path.join(IMG_DIR, filename)
    iw, ih = PILImage.open(path).size
    src_ar = iw / ih
    box_ar = w / h
    pic = slide.shapes.add_picture(path, int(l), int(t), width=int(w), height=int(h))
    if src_ar > box_ar:          # source too wide -> crop left/right
        keep = box_ar / src_ar
        c = (1 - keep) / 2
        pic.crop_left = c
        pic.crop_right = c
    elif src_ar < box_ar:        # source too tall -> crop top/bottom
        keep = src_ar / box_ar
        c = (1 - keep) / 2
        pic.crop_top = c
        pic.crop_bottom = c
    return pic


def framed_photo(slide, l, t, w, img_h, filename, caption):
    """White frame + cover-cropped photo + navy caption strip. Returns total height."""
    pad = Inches(0.08)
    cap_h = Inches(0.5)
    total_h = img_h + pad * 2 + cap_h
    box(slide, l, t, w, total_h, WHITE, line=BLUE, lw=1.5)
    cover_pic(slide, l + pad, t + pad, w - pad * 2, img_h, filename)
    cap_t = t + pad + img_h + Inches(0.02)
    box(slide, l + pad, cap_t, w - pad * 2, cap_h, NAVY)
    txt(slide, l + pad, cap_t, w - pad * 2, cap_h, caption, F_CAP, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return total_h


# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
PANEL_W = Inches(6.25)
cover_pic(s, PANEL_W, 0, SW - PANEL_W, SH, "patna_cover.jpg")
box(s, 0, 0, PANEL_W, SH, NAVY)
box(s, Inches(0.65), Inches(2.3), Inches(0.16), Inches(2.35), ORANGE)
txt(s, Inches(0.95), Inches(1.25), Inches(5.1), Inches(1.5),
    "Parking & Traffic Management", 32, WHITE, bold=True)
txt(s, Inches(0.95), Inches(2.5), Inches(5.1), Inches(0.8),
    "Issues and Challenges", 23, ORANGE, bold=True)
txt(s, Inches(0.95), Inches(3.45), Inches(5.1), Inches(0.6),
    "A Session for City Managers  •  Patna", 16, WHITE)
txt(s, Inches(0.95), Inches(5.45), Inches(5.1), Inches(0.45),
    "Session by", 13, PALE)
txt(s, Inches(0.95), Inches(5.8), Inches(5.1), Inches(0.55),
    "Pourush Agarwal", 21, WHITE, bold=True)
txt(s, Inches(0.95), Inches(6.35), Inches(5.1), Inches(0.5),
    "DySP Traffic Patna", 16, ORANGE, bold=True)

# =====================================================================
# SLIDE 2 — Agenda
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Session Agenda", 2)
bullets(s, Inches(0.8), Inches(1.5), Inches(11.6), Inches(5.6), [
    "Why parking & traffic matter for a city",
    "Key parking issues & the role of sign boards",
    "Traffic management challenges",
    "Special situations: VIP movement & metro diversion",
    "Root causes and their impact",
    "Solutions — smarter parking, signage & real-time updates",
    "Smart technology in traffic management",
    "Case study: Pune parking reform",
    "Action roadmap for Patna",
], gap=13)

# =====================================================================
# SLIDE 3 — Why it matters
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Why Parking & Traffic Matter", 3)
bullets(s, Inches(0.8), Inches(1.5), Inches(11.6), Inches(5.6), [
    "Traffic reflects a city's quality of life & governance",
    "Poor parking eats into scarce public road space",
    "Congestion means lost productivity, fuel and time",
    "It directly affects safety, emergency access & emissions",
    "Well-managed mobility builds public trust & attracts investment",
    "Orderly streets improve pedestrian safety and business footfall",
], gap=16)

# =====================================================================
# SLIDE 4 — Key Parking Issues (2 photos)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Key Parking Issues in Patna", 4)
bullets(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(2.9), [
    "Illegal roadside parking on both sides of arterial roads",
    "Haphazard parking chokes market roads, as seen in Wakarganj",
    "Vehicles ignore signage — parking right under 'No Parking' boards",
    "Autos & e-rickshaws halt anywhere near stations & markets",
    "Very little organised, paid off-street parking",
], gap=7)
# two photos, side by side
PH_W = Inches(4.75)
PH_IMG_H = Inches(2.28)
gap = Inches(0.35)
start_x = (SW - (PH_W * 2 + gap)) / 2
framed_photo(s, start_x, Inches(4.2), PH_W, PH_IMG_H, "patna_noparking.jpg",
             "Parking right under a 'No Parking' board")
framed_photo(s, start_x + PH_W + gap, Inches(4.2), PH_W, PH_IMG_H,
             "patna_haphazard_wakarganj.jpg", "Haphazard parking in Wakarganj")

# =====================================================================
# SLIDE 5 — Traffic Management Challenges (bigger photo)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Traffic Management Challenges", 5)
bullets(s, Inches(0.6), Inches(1.45), Inches(6.85), Inches(5.6), [
    "Mixed traffic — cars, 2-wheelers, autos, e-rickshaws, carts",
    "Traffic peaks during school and office hours at key locations",
    "Severe manpower shortage",
    "Weak public transport pushes people to private vehicles",
    "Only ~8% of city area is road (well below the 25% norm)",
    "Frequent VIP movement disrupts normal traffic flow",
    "Metro construction forces route diversions & lane closures",
    "Encroachment & construction shrink usable road width",
], gap=11)
framed_photo(s, Inches(7.65), Inches(1.5), Inches(5.15), Inches(4.55),
             "patna_regulation.jpg", "Patna Traffic Police regulating a busy junction")

# =====================================================================
# SLIDE 6 — Special Traffic Situations (smaller boxes, wider photo)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Special Traffic Situations", 6)
BOX_W = Inches(5.55)
BOX_H = Inches(1.85)
bx1 = Inches(0.75)
bx2 = Inches(7.05)
box(s, bx1, Inches(1.4), BOX_W, BOX_H, WHITE, line=PURPLE, lw=2)
txt(s, bx1 + Inches(0.22), Inches(1.5), BOX_W - Inches(0.4), Inches(0.45),
    "VIP Movement", F_CARD_H, PURPLE, bold=True)
bullets(s, bx1 + Inches(0.22), Inches(2.0), BOX_W - Inches(0.4), Inches(1.2), [
    "Route clearance & holding of traffic during transit",
    "Advance planning with escort & pilot coordination",
], size=F_CARD_T, gap=5)
box(s, bx2, Inches(1.4), BOX_W, BOX_H, WHITE, line=BLUE, lw=2)
txt(s, bx2 + Inches(0.22), Inches(1.5), BOX_W - Inches(0.4), Inches(0.45),
    "Metro Diversion", F_CARD_H, BLUE, bold=True)
bullets(s, bx2 + Inches(0.22), Inches(2.0), BOX_W - Inches(0.4), Inches(1.2), [
    "Lane closures around metro construction sites",
    "Clear signage + advance public communication",
], size=F_CARD_T, gap=5)
WIDE_W = Inches(8.6)
framed_photo(s, (SW - WIDE_W) / 2, Inches(3.5), WIDE_W, Inches(2.9),
             "patna_boringroad.jpg",
             "On-ground regulation at Boring Road crossing, Patna")

# =====================================================================
# SLIDE 7 — Root Causes
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Root Causes Behind the Chaos", 7)
causes = [("Rapid Motorisation", "Vehicle ownership rising faster than infrastructure"),
          ("Weak Enforcement", "Low penalties and poor signage compliance"),
          ("No Parking Policy", "Parking treated as free & unlimited"),
          ("Land-use Mismatch", "Dense commercial zones with no parking norms"),
          ("Weak Public Transport", "Few reliable last-mile options"),
          ("Fragmented Governance", "Multiple agencies, unclear ownership")]
for i, (h, d) in enumerate(causes):
    cx = Inches(0.7) + Inches(4.05) * (i % 3)
    cy = Inches(1.55) + Inches(2.5) * (i // 3)
    box(s, cx, cy, Inches(3.8), Inches(2.25), WHITE, line=BLUE)
    txt(s, cx + Inches(0.22), cy + Inches(0.2), Inches(3.4), Inches(0.6),
        h, F_CARD_H, BLUE, bold=True)
    txt(s, cx + Inches(0.22), cy + Inches(0.85), Inches(3.4), Inches(1.2),
        d, F_CARD_T, GREY)

# =====================================================================
# SLIDE 8 — Impact (smaller boxes + big photo)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Impact of Poor Management", 8)
cols = [("Economic", ["Lost man-hours in jams", "Higher fuel & logistics cost",
                      "Reduced business footfall"], ORANGE),
        ("Environmental", ["Higher vehicular emissions", "Idling pollution & noise",
                           "Poor air quality (AQI)"], GREEN),
        ("Social & Safety", ["More accidents & fatalities", "Blocked emergency vehicles",
                             "Stress & road rage"], RED)]
CW = Inches(2.45)
cx = Inches(0.6)
for title, items, col in cols:
    box(s, cx, Inches(1.5), CW, Inches(0.6), col)
    txt(s, cx, Inches(1.5), CW, Inches(0.6), title, F_CARD_T + 1, WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, cx, Inches(2.1), CW, Inches(3.1), WHITE, line=col)
    bullets(s, cx + Inches(0.15), Inches(2.28), CW - Inches(0.3), Inches(2.8),
            items, size=F_CARD_T, gap=10)
    cx += CW + Inches(0.12)
framed_photo(s, Inches(8.55), Inches(1.5), Inches(4.25), Inches(3.15),
             "patna_impact.jpg", "Congestion on a Patna arterial road")

# =====================================================================
# SLIDE 9 — Smarter Parking (bigger photo)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Solutions: Smarter Parking", 9)
bullets(s, Inches(0.6), Inches(1.5), Inches(6.6), Inches(5.5), [
    "Priced, demand-based on-street parking (pay-and-park)",
    "Clearly marked bays with barricades & discipline",
    "Multi-level & organised off-street parking",
    "Designated auto and e-rickshaw stands near hubs",
    "Strict enforcement: towing, fines and CCTV",
    "Parking norms in building bye-laws for new projects",
], gap=15)
framed_photo(s, Inches(7.55), Inches(1.5), Inches(5.25), Inches(4.55),
             "patna_properparking.jpg", "Organised auto stand with barricades — Patna")

# =====================================================================
# SLIDE 10 — Traffic Management Solutions (two columns)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Solutions: Traffic Management", 10)
bullets(s, Inches(0.6), Inches(1.65), Inches(6.2), Inches(5.6), [
    "Clear, well-placed sign boards with strict enforcement",
    "Real-time updates to public — social media, apps & alerts",
    "Advance communication for VIP movement & metro diversions",
    "Adaptive / synchronised signals on key corridors (ITS)",
    "Dedicated lanes for buses & non-motorised transport",
    "Community awareness by NCC",
], gap=26)
bullets(s, Inches(6.95), Inches(1.65), Inches(5.9), Inches(5.6), [
    "Joint and coordinated anti-encroachment drives",
    "Elevated corridors and flyovers",
    ("Karbigahiya–Mithapur Flyover", True),
    ("Mithapur–Sipara Elevated Road", True),
    "Ring road development",
    "Smart management : ICCC",
    "Strengthen public transport & last-mile connectivity",
], gap=22)

# =====================================================================
# SLIDE 11 — Smart Technology
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "The Role of Smart Technology", 11)
tech = [("Smart Sensors", "Detect free and occupied parking bays in real time"),
        ("Digital Sign Boards", "Dynamic messages for diversions & warnings"),
        ("ANPR Cameras", "Automatic number-plate based enforcement"),
        ("Adaptive Signals", "Signal timing based on live traffic flow"),
        ("Real-time Updates", "Live traffic & diversion information to citizens"),
        ("Cashless Payment", "Transparent, leak-proof parking revenue")]
for i, (h, d) in enumerate(tech):
    cx = Inches(0.7) + Inches(4.05) * (i % 3)
    cy = Inches(1.55) + Inches(2.5) * (i // 3)
    box(s, cx, cy, Inches(3.8), Inches(2.25), NAVY)
    txt(s, cx + Inches(0.22), cy + Inches(0.2), Inches(3.4), Inches(0.6),
        h, F_CARD_H, ORANGE, bold=True)
    txt(s, cx + Inches(0.22), cy + Inches(0.85), Inches(3.4), Inches(1.2),
        d, F_CARD_T, WHITE)

# =====================================================================
# SLIDE 12 — Case Study: Pune (+ photo)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Case Study: Pune Parking Reform", 12)
box(s, Inches(0.6), Inches(1.5), Inches(0.14), Inches(5.0), ORANGE)
txt(s, Inches(0.9), Inches(1.45), Inches(6.6), Inches(0.55),
    "Pune's on-street parking management (with ITDP India)",
    F_CARD_H, BLUE, bold=True)
bullets(s, Inches(0.9), Inches(2.1), Inches(6.6), Inches(4.6), [
    "City parking policy adopted — parking priced, not free",
    "Demand-based pricing in busy zones & peak hours",
    "New no-parking zones",
    "Metro–parking integration — app-based booking and discount for ticket holders",
    "Cashless, self-paying system for transparent revenue",
    "Result: less haphazard parking & smoother flow",
], gap=12)
framed_photo(s, Inches(7.9), Inches(1.5), Inches(4.9), Inches(3.2),
             "patna_case_photo.png", "Integrated parking & transit facility")

# =====================================================================
# SLIDE 13 — Action Roadmap
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT)
header(s, "Action Roadmap for Patna", 13)
phases = [("Short Term (0-6 months)",
           ["Fix hotspots — Boring Road, Station Road", "Demarcate bays with clear signage",
            "Clear encroachments", "Enforcement drive and towing"], ORANGE),
          ("Medium Term (6-18 months)",
           ["Adopt a city parking policy", "Pay-and-park with mobile app",
            "Real-time updates channel", "Metro diversion planning"], BLUE),
          ("Long Term (18 months +)",
           ["Multi-level parking", "City-wide ITS and CCTV",
            "Integrated mobility plan", "Data-driven governance"], GREEN)]
cx = Inches(0.7)
for title, items, col in phases:
    box(s, cx, Inches(1.5), Inches(3.95), Inches(0.7), col)
    txt(s, cx, Inches(1.5), Inches(3.95), Inches(0.7), title, F_CARD_T + 1, WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, cx, Inches(2.2), Inches(3.95), Inches(3.4), WHITE, line=col)
    bullets(s, cx + Inches(0.2), Inches(2.4), Inches(3.6), Inches(3.1),
            items, size=F_CARD_T, gap=13)
    cx += Inches(4.15)
box(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.85), NAVY)
txt(s, Inches(0.95), Inches(5.85), Inches(11.4), Inches(0.85),
    "Already begun: In late 2025 Patna traffic police banned car and auto parking on "
    "Income Tax Golambar–Dak Bungalow and Budh Marg–GPO Golambar.",
    F_CARD_T, WHITE, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# SLIDE 14 — Key Takeaways
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, Inches(1.15), SW, Inches(0.08), ORANGE)
txt(s, Inches(0.5), Inches(0.25), Inches(11.4), Inches(0.8), "Key Takeaways",
    F_TITLE, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.85), Inches(1.5), Inches(11.6), Inches(5.6), [
    "Parking is a policy problem, not just an infrastructure one",
    "Sign boards must be backed by real enforcement",
    "Plan ahead for VIP movement and metro diversions",
    "Real-time updates keep citizens informed and reduce chaos",
    "Joint and coordinated efforts by all stakeholders",
    "Holistic approach by city managers",
    "Patna can adapt proven models like Pune's reform",
], color=WHITE, gap=13)

# =====================================================================
# SLIDE 15 — Thank You
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
THANK_W = Inches(6.1)
cover_pic(s, SW - THANK_W, 0, THANK_W, SH, "patna_thanks.jpg")
box(s, 0, 0, SW - THANK_W, SH, NAVY)
box(s, Inches(0.8), Inches(2.55), Inches(0.16), Inches(1.9), ORANGE)
txt(s, Inches(1.15), Inches(2.35), Inches(5.4), Inches(1.0), "Thank You", 44,
    WHITE, bold=True)
txt(s, Inches(1.15), Inches(3.45), Inches(5.4), Inches(0.6),
    "Questions & Discussion", 21, ORANGE, bold=True)
txt(s, Inches(1.15), Inches(4.75), Inches(5.4), Inches(0.45),
    "Session by", 13, PALE)
txt(s, Inches(1.15), Inches(5.1), Inches(5.4), Inches(0.5),
    "Pourush Agarwal", 20, WHITE, bold=True)
txt(s, Inches(1.15), Inches(5.6), Inches(5.4), Inches(0.5),
    "DySP Traffic Patna", 16, ORANGE, bold=True)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Parking_Traffic_Management_Patna_v2.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
