#!/usr/bin/env python3
"""Fix hanging-indent alignment for every bullet paragraph across all slides.

Problem: bullets are typed as literal text ("•  text" / "     – text"),
so when a line wraps, the second line starts at the textbox's left edge
instead of lining up under the first word after the bullet.

Fix: set a proper hanging indent (marL = left margin for wrapped lines,
indent = negative first-line offset so the bullet/dash sits back at the
margin) on every paragraph that starts with "•" or a sub-bullet "–".
Content/text is NOT changed — only paragraph indentation.
"""
from pptx import Presentation
from pptx.util import Emu, Inches
import os

BASE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(BASE, "Parking_Traffic_Management_Patna_v2.pptx")

MAIN_MARL = Inches(0.30)     # where wrapped text of a "•" bullet lines up
MAIN_INDENT = -Inches(0.30)  # pulls the "•" back to the textbox edge

SUB_MARL = Inches(0.62)      # where wrapped text of a "  – " sub-bullet lines up
SUB_INDENT = -Inches(0.32)   # pulls the "–" back under the parent bullet's text

prs = Presentation(PPTX)

fixed = 0
scanned = 0
for s_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            text = p.text
            if not text:
                continue
            scanned += 1
            pPr = p._p.get_or_add_pPr()
            if text.startswith("•"):
                pPr.set("marL", str(int(MAIN_MARL)))
                pPr.set("indent", str(int(MAIN_INDENT)))
                fixed += 1
            elif text.lstrip().startswith("–"):
                pPr.set("marL", str(int(SUB_MARL)))
                pPr.set("indent", str(int(SUB_INDENT)))
                fixed += 1

prs.save(PPTX)
print(f"Scanned {scanned} paragraphs, applied hanging indent to {fixed} bullet paragraphs.")
print("Saved:", PPTX)
