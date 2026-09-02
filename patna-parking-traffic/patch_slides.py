#!/usr/bin/env python3
"""Patch the existing pptx in place:
   - Slide 9: replace the case-study photo with pune_parking_reform.png
   - Slide 16: replace the thank-you photo with the enhanced (sharper/brighter) thanks image
Everything else (including the user's manual edit to Slide 1) is left untouched.
"""
from pptx import Presentation
from pptx.util import Inches
import os

BASE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(BASE, "Parking_Traffic_Management_Patna_v2.pptx")
IMG_DIR = os.path.join(BASE, "ppt_images")

prs = Presentation(PPTX)
slides = prs.slides

def replace_picture(slide, new_image_path, crop_mode="cover"):
    """Find the picture shape on the slide, remember its box, delete it,
    and insert the new image in the same box (cover-crop, no stretch)."""
    from PIL import Image as PILImage
    pic_shape = None
    for shp in slide.shapes:
        if shp.shape_type == 13:  # PICTURE
            pic_shape = shp
            break
    if pic_shape is None:
        raise RuntimeError("No picture shape found on this slide")

    l, t, w, h = pic_shape.left, pic_shape.top, pic_shape.width, pic_shape.height
    # remove old picture
    pic_shape._element.getparent().remove(pic_shape._element)

    iw, ih = PILImage.open(new_image_path).size
    src_ar = iw / ih
    box_ar = w / h
    new_pic = slide.shapes.add_picture(new_image_path, l, t, width=w, height=h)
    if crop_mode == "cover":
        if src_ar > box_ar:
            keep = box_ar / src_ar
            c = (1 - keep) / 2
            new_pic.crop_left = c
            new_pic.crop_right = c
        elif src_ar < box_ar:
            keep = src_ar / box_ar
            c = (1 - keep) / 2
            new_pic.crop_top = c
            new_pic.crop_bottom = c
    return new_pic


# ---- FIX previous mistake: restore Slide 9 (index 8, "Smarter Parking") ----
slide9 = slides[8]
replace_picture(slide9, os.path.join(IMG_DIR, "patna_properparking.jpg"))
print("Slide 9 (Smarter Parking) photo restored to patna_properparking.jpg")

# ---- Case Study: Pune Parking Reform slide (index 11) -> pune_parking_reform.png ----
slide_case_study = slides[11]
replace_picture(slide_case_study, os.path.join(IMG_DIR, "pune_parking_reform.png"))
print("Case Study slide (index 11) photo replaced with pune_parking_reform.png")

# ---- Slide 16 (index 15): Thank-you photo -> enhanced thanks image ----
# Note: presentation currently has 15 slides in this file (0-14); the
# thank-you slide is the LAST slide regardless of its printed number.
last_idx = len(slides) - 1
slide_last = slides[last_idx]
replace_picture(slide_last, os.path.join(IMG_DIR, "patna_thanks.jpg"))
print(f"Last slide (index {last_idx}) photo replaced with enhanced thanks image")

prs.save(PPTX)
print("Saved:", PPTX)
