#!/usr/bin/env python3
"""Re-layout the title slide's left-panel text so it flows naturally:
Title -> Subtitle -> (small gap) -> Session by / Name / Role.
No content changes, only vertical positions."""
from pptx import Presentation
from pptx.util import Inches
import os

BASE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(BASE, "Parking_Traffic_Management_Patna_v2.pptx")

prs = Presentation(PPTX)
slide = prs.slides[0]

positions = {
    "Parking & Traffic Management": Inches(1.25),
    "Issues and Challenges": Inches(2.55),
    "Session by": Inches(3.55),
    "Pourush Agarwal": Inches(3.9),
    "DySP Traffic Patna": Inches(4.45),
}

for shp in slide.shapes:
    if shp.has_text_frame:
        text = shp.text_frame.text.strip()
        if text in positions:
            shp.top = positions[text]

# Re-center the orange accent bar next to the text block
for shp in slide.shapes:
    if shp.shape_type == 1 and shp.width < Inches(0.3) and shp.height > Inches(1):
        shp.top = Inches(1.25)
        shp.height = Inches(3.7)

prs.save(PPTX)
print("Saved:", PPTX)
